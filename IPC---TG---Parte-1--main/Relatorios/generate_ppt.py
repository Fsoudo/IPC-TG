import sys
import subprocess
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt

prs = Presentation()

# Helper to add slides
def add_slide(title, bullet_points):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    title_placeholder.text = title
    tf = body_shape.text_frame
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Simple sub-bullet logic
        if point.startswith("  - "):
            p.text = point[4:]
            p.level = 1
        else:
            p.text = point
            p.level = 0

# Slide 1
slide_layout = prs.slide_layouts[0] # Title Slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Apresentação do Projeto: EasyMed"
slide.shapes.placeholders[1].text = "Interação Pessoa-Computador · TG1\nEquipa 14: Francisco Soudo e Miguel Pauzinho\nAno Letivo: 2025/2026 · IPBeja / ESTIG"

# Slide 2
add_slide("Análise de Sistemas Semelhantes", [
    "MyTherapy (Francisco): Design limpo e foco em lembretes, mas com risco de sobrecarga funcional devido a configurações complexas.",
    "Medisafe (Miguel): Foco na gestão de medicação, avaliada pelo ponto de vista de acessibilidade e usabilidade para o público-alvo principal.",
    "Conclusão de Equipa: A EasyMed deve focar-se na simplicidade e na eliminação de funcionalidades excessivas, priorizando a acessibilidade."
])

# Slide 3
add_slide("Personas", [
    "Persona 1: António Ferreira (Miguel)",
    "  - 72 anos, baixa literacia digital (Utilizador Principal).",
    "  - Precisa de alertas simples e botões grandes para confirmar as tomas.",
    "Persona 2: Carla Sousa (Francisco)",
    "  - 46 anos, enfermeira, cuidadora familiar (Utilizador Secundário).",
    "  - Precisa de monitorizar o pai à distância e receber alertas rápidos em caso de esquecimento."
])

# Slide 4
add_slide("Cenários de Utilização", [
    "Cenário 1: Confirmar Toma de Medicamento (António)",
    "  - Recebe alerta sonoro, vai buscar o medicamento e confirma com um botão na notificação expandida.",
    "Cenário 2: Receber Notificação de Toma Ignorada (Carla)",
    "  - Recebe alerta crítico, verifica que o pai não tomou (código vermelho) e liga-lhe imediatamente pelo atalho na app."
])

# Slide 5
add_slide("Análise de Tarefas (HTA)", [
    "HTA Cenário 1: Decomposição em 4 passos simples focados na ação imediata de tomar e confirmar.",
    "HTA Cenário 2: Decomposição focada na reação da cuidadora (receber alerta, consultar detalhe, contactar pai, verificar).",
    "Abordagem da Equipa: Foco em reduzir ao máximo o número de toques (cliques) para atingir o objetivo em ambos os cenários."
])

# Slide 6
add_slide("Estilos e Dispositivos de Interação", [
    "Estilos: Manipulação direta (toques largos/swipe), Menus/Listas simples e Notificações proativas (push).",
    "Dispositivos: Ecrã tátil de smartphone Android.",
    "Acessibilidade: Botões de grandes dimensões (≥ 48x48dp), uso complementar de som e vibração para alertas de toma."
])

# Slide 7
add_slide("Princípios de Usabilidade (Nielsen)", [
    "Visibilidade do Estado do Sistema: Códigos de cores claros (Verde = Tomado, Vermelho = Ignorado) e feedback imediato.",
    "Prevenção de Erros: Exigência de swipe ou botões de confirmação claros para evitar toques acidentais (fat-finger errors).",
    "Correspondência com o Mundo Real: Termos não técnicos e listas cronológicas semelhantes a uma agenda diária.",
    "Estética e Design Minimalista: Ecrãs limpos sem distrações, essenciais tanto para o idoso como para a urgência da cuidadora."
])

# Slide 8
add_slide("Protótipos Figma e Navegação", [
    "Cenário 1 (António): Notificação -> Notificação Expandida -> Confirmação -> Ecrã Principal.",
    "Cenário 2 (Carla): Notificação -> Dashboard do Cuidador -> Detalhe das Tomas -> Histórico Semanal.",
    "Guia de Estilo Unificado: Cor primária #4A90E2, tipografia Roboto e coerência visual entre ambas as perspetivas."
])

# Slide 9
add_slide("Organização do Trabalho e Scrum", [
    "Metodologia: Divisão de tarefas estruturada em Trello com base no framework Scrum.",
    "Trabalho Colaborativo: Definição conjunta de interações e estilos, com cenários distintos que se complementam.",
    "Ferramentas: Comunicação via WhatsApp, versionamento de relatórios em GitHub e design colaborativo em Figma."
])

# Slide 10
add_slide("Utilização de IA no Projeto", [
    "Aplicação Ética e Produtiva: O Claude (Anthropic) foi usado como co-piloto por ambos os elementos da equipa.",
    "Casos de Uso: Estruturação de textos, brainstorming para personas, geração de código Mermaid e suporte ao layout no Figma.",
    "Validação: Todo o conteúdo gerado foi validado, revisto e aprovado manualmente, garantindo alinhamento com a teoria IPC."
])

# Slide 11
add_slide("Reflexão Final da Equipa", [
    "Acessibilidade como Prioridade (Miguel): O design tem de se adaptar ao utilizador com literacia reduzida através de feedback imediato.",
    "O Papel do Utilizador Secundário (Francisco): A app só funciona se a cuidadora tiver uma interface rápida e sem sobrecarga em crise.",
    "Conclusão: O Figma permitiu materializar estes conceitos num protótipo visualmente coerente e próximo de uma aplicação real."
])

output_file = os.path.join(os.path.dirname(__file__), "Apresentacao_Global_EasyMed.pptx")
prs.save(output_file)
print(f"Presentation generated successfully at: {output_file}")
